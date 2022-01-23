from pptx import Presentation

class AutoMaker:

    def __init__(self):
        self.ppt = Presentation()
        self.copy_slide = []    

    def CopyAndPaste(self, *args):
        for i in args:
            name = "".join(i)
            source_slide = Presentation(name+'.pptx')
            for x in range(10):
                try:
                    slide_layout = source_slide.slide_layouts[x]
                    
                except:
                    break
                
                
              #  self.copy_slide = self.ppt.slides.add_slide(slide_layout)
                self.ppt.slides.add_slide(slide_layout)
        
    

    def SaveAsPPTX(self, name):
        self.ppt.save(name + '.pptx')



        
